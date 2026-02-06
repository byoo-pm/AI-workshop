
import os
import sys
from pptx import Presentation

# Force utf-8 output
sys.stdout.reconfigure(encoding='utf-8')

def read_pptx(file_path):
    try:
        prs = Presentation(file_path)
        full_text = []
        for i, slide in enumerate(prs.slides):
            slide_text = []
            slide_title = "No Title"
            if slide.shapes.title:
                slide_title = slide.shapes.title.text
            
            slide_text.append(f"--- Slide {i+1}: {slide_title} ---")
            
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                     if shape.text.strip():
                        slide_text.append(shape.text)
                
                # Extract text from tables
                if shape.has_table:
                    for row in shape.table.rows:
                        row_text = []
                        for cell in row.cells:
                            if cell.text.strip():
                                row_text.append(cell.text.strip())
                        if row_text:
                            slide_text.append(" | ".join(row_text))

            full_text.append("\n".join(slide_text))
        return '\n\n'.join(full_text)
    except Exception as e:
        return f"Error reading {file_path}: {str(e)}"

file_path = r"c:\Users\byulh\Documents\Antigravity\career_advancement\ai_survey\assets\intel\NCE Prod. Marketing team meeting - 06.02.2026.pptx"

print("--- START OF PRESENTATION ---")
content = read_pptx(file_path)
with open("ppt_content.txt", "w", encoding="utf-8") as f:
    f.write(content)
print("--- END OF PRESENTATION --- Content written to ppt_content.txt")
